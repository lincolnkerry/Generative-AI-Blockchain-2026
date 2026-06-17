#!/usr/bin/env python3
"""Compile slide PNGs into PDF and PPTX."""
import sys
import glob
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

def compile_pdf(slide_dir, output_pdf):
    pngs = sorted(glob.glob(os.path.join(slide_dir, 'slide-*.png')))
    if not pngs:
        print(f"No slides found in {slide_dir}")
    imgs = [Image.open(p).convert('RGB') for p in pngs]
    # Use first image's save with optimize=False to avoid JPEG encoder issue
    imgs[0].save(output_pdf, 'PDF', save_all=True, append_images=imgs[1:])

def compile_pptx(slide_dir, output_pptx):
    pngs = sorted(glob.glob(os.path.join(slide_dir, 'slide-*.png')))
    if not pngs:
        print(f"No slides found in {slide_dir}")
        return
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for png in pngs:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        slide.shapes.add_picture(png, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    prs.save(output_pptx)
    print(f"PPTX: {output_pptx} ({len(pngs)} slides)")

if __name__ == '__main__':
    for lang in ['en', 'kr']:
        slide_dir = f'slides/slides_{lang}'
        compile_pdf(slide_dir, f'slides/presentation_{lang}.pdf')
        compile_pptx(slide_dir, f'slides/presentation_{lang}.pptx')
