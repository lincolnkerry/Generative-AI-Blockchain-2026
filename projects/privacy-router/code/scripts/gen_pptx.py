"""Generate Privacy Router presentation as PPTX (Bold Signal style)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Design tokens (Bold Signal) ─────────────────────────────────────
BG = RGBColor(0x1A, 0x1A, 0x1A)
BG2 = RGBColor(0x2D, 0x2D, 0x2D)
CARD = RGBColor(0xFF, 0x57, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x9C, 0xA3, 0xAF)
WARM = RGBColor(0xFF, 0x8A, 0x65)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
RED = RGBColor(0xF4, 0x43, 0x36)
CYAN = RGBColor(0x00, 0xBC, 0xD4)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
LIGHT_BG = RGBColor(0x22, 0x22, 0x22)

W = Inches(13.333)  # 1920px @ 96dpi → 20"
H = Inches(7.5)     # 1080px → 15" — but we use standard 16:9

prs = Presentation()
prs.slide_width = Emu(12192000)   # 13.333"
prs.slide_height = Emu(6858000)   # 7.5"
SW = prs.slide_width
SH = prs.slide_height


def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=18,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
             font_name='Arial'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rich_text(slide, left, top, width, height, runs, alignment=PP_ALIGN.LEFT):
    """runs: list of (text, size, color, bold, font_name)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    for i, (text, size, color, bold, font_name) in enumerate(runs):
        if i == 0:
            run = p.runs[0] if p.runs else p.add_run()
            run.text = text
        else:
            run = p.add_run()
            run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = font_name or 'Arial'
    return txBox


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, label, title, body,
             fill=LIGHT_BG, border=None):
    rect = add_rect(slide, left, top, width, height, fill, border)
    rect.adjustments[0] = 0.04  # corner radius
    add_text(slide, left + Inches(0.3), top + Inches(0.2), width - Inches(0.6), Inches(0.3),
             label, 10, WARM, True)
    add_text(slide, left + Inches(0.3), top + Inches(0.5), width - Inches(0.6), Inches(0.5),
             title, 16, WHITE, True)
    add_text(slide, left + Inches(0.3), top + Inches(1.0), width - Inches(0.6), height - Inches(1.3),
             body, 12, GRAY)


def footer(slide, left_text, right_text):
    add_text(slide, Inches(0.8), SH - Inches(0.55), Inches(5), Inches(0.3),
             left_text, 10, GRAY)
    add_text(slide, SW - Inches(3.8), SH - Inches(0.55), Inches(3), Inches(0.3),
             right_text, 10, GRAY, alignment=PP_ALIGN.RIGHT)


def section_num(slide, num):
    add_text(slide, Inches(0.4), Inches(0.2), Inches(3), Inches(1.5),
             f"{num:02d}", 80, RGBColor(0x2D, 0x2D, 0x2D), font_name='Arial')


def nav_breadcrumb(slide, items, active_idx):
    """Draw nav breadcrumb top-right."""
    parts = []
    for i, item in enumerate(items):
        if i > 0:
            parts.append((" / ", 10, RGBColor(0x55, 0x55, 0x55), False, "Arial"))
        if i == active_idx:
            parts.append((item, 10, WARM, True, "Arial"))
        else:
            parts.append((item, 10, RGBColor(0x66, 0x66, 0x66), False, "Arial"))
    add_rich_text(slide, SW - Inches(5), Inches(0.35), Inches(4.2), Inches(0.3),
                  parts, PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(s)
section_num(s, 1)
add_text(s, Inches(0.8), SH - Inches(2.8), Inches(6), Inches(1),
         "Privacy Router", 52, WHITE, True)
add_text(s, Inches(0.8), SH - Inches(2.0), Inches(6), Inches(0.6),
         "Sovereignty-First Butler for Korean Researchers", 18, GRAY)

# Accent card
card_l = SW - Inches(5.5)
card_t = Inches(1.8)
card_w = Inches(4.7)
card_h = Inches(4.2)
add_rect(s, card_l, card_t, card_w, card_h, CARD)
add_text(s, card_l + Inches(0.4), card_t + Inches(0.3), Inches(3.5), Inches(0.3),
         "HOW IT WORKS", 10, DARK_TEXT, True)
add_text(s, card_l + Inches(0.4), card_t + Inches(0.7), Inches(3.8), Inches(0.8),
         "Extractor → Judge → Router", 22, DARK_TEXT, True)
add_text(s, card_l + Inches(0.4), card_t + Inches(1.5), Inches(3.8), Inches(1.5),
         "Detects sensitive spans in Korean and English prompts, decides whether to mask, block, or process locally — then routes to the right model.",
         13, DARK_TEXT)
# Chips
for i, (label, bg_c) in enumerate([("Local (Gemma 4)", RGBColor(0x33, 0x33, 0x33)),
                                     ("Cloud (Gemini)", RGBColor(0xFF, 0xFF, 0xFF)),
                                     ("Mask & Send", RGBColor(0x33, 0x33, 0x33))]):
    chip_x = card_l + Inches(0.4) + Inches(i * 1.4)
    chip_y = card_t + Inches(3.3)
    r = add_rect(s, chip_x, chip_y, Inches(1.3), Inches(0.35), bg_c)
    tc = DARK_TEXT if bg_c == WHITE else DARK_TEXT
    add_text(s, chip_x + Inches(0.05), chip_y + Inches(0.03), Inches(1.2), Inches(0.3),
             label, 9, tc, True, PP_ALIGN.CENTER)

footer(s, "GenAI & Blockchain · Term Project · June 2026", "DH. Kim & M. Saadati")


# ════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
section_num(s, 2)
nav_breadcrumb(s, ["Title", "Problem", "Architecture", "Eval"], 1)
add_text(s, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
         "The Problem", 40, WHITE, True)
add_text(s, Inches(0.8), Inches(1.1), Inches(8), Inches(0.4),
         "Korean researchers paste sensitive data into LLMs without protection", 16, GRAY)

cards = [
    ("REAL RISK", "Resident Registration Numbers",
     "Korean RRN (주민등록번호) is a 13-digit national ID. Existing guardrails miss it — regex doesn't understand context."),
    ("REAL RISK", "Unpublished Research",
     "Research ideas, lab budgets, submission plans — competitive intelligence that regex can't detect."),
    ("GAP", "No Korean Privacy Layer",
     "Cloudflare LLM Firewall, AWS Bedrock Guardrails — all English-centric. None understand Korean PII."),
    ("OUR ANSWER", "Privacy Router",
     "Korean-first prompt firewall that detects, masks, and routes — keeping sensitive data local."),
]
for i, (label, title, body) in enumerate(cards):
    col = i % 2
    row = i // 2
    cl = Inches(0.8) + col * Inches(6.1)
    ct = Inches(1.8) + row * Inches(2.5)
    cw = Inches(5.7)
    ch = Inches(2.2)
    border = CARD if i == 3 else RGBColor(0x33, 0x33, 0x33)
    fill = RGBColor(0x26, 0x20, 0x1A) if i == 3 else LIGHT_BG
    add_card(s, cl, ct, cw, ch, label, title, body, fill, border)

footer(s, "GenAI & Blockchain · Term Project", "2 / 12")


# ════════════════════════════════════════════════════════════════════
# SLIDE 3 — Three-Harm Test
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
section_num(s, 3)
nav_breadcrumb(s, ["Title", "Problem", "Framework", "Eval"], 2)
add_text(s, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
         "Three-Harm Test", 40, WHITE, True)
add_text(s, Inches(0.8), Inches(1.1), Inches(8), Inches(0.4),
         "Every sentence is evaluated against three dimensions of harm", 16, GRAY)

harms = [
    ("IDENTITY", RED, "Can this identify a person?",
     "RRN, phone, email, passport, name+institution pairs\n\n\"901212-1234567\" → RRN\n\"GIST 김동현\" → STUDENT_AFFILIATION"),
    ("COMPETITIVE", RGBColor(0xFF, 0x98, 0x00), "Can a rival benefit?",
     "Unpublished research, internal decisions, strategy\n\n\"TSMC 3nm 공정 채택\"\n→ FABRICATION_DECISION"),
    ("SAFETY", CYAN, "Can exposure cause harm?",
     "Credentials, internal URLs, budgets, salaries\n\n\"@newpurk2025\" → PASSWORD\n\"예산 5억원\" → BUDGET"),
]
for i, (tag, color, title, body) in enumerate(harms):
    cl = Inches(0.8) + i * Inches(4.1)
    ct = Inches(1.8)
    cw = Inches(3.8)
    ch = Inches(4.8)
    rect = add_rect(s, cl, ct, cw, ch, LIGHT_BG, RGBColor(0x33, 0x33, 0x33))
    add_rect(s, cl, ct, cw, Inches(0.06), color)
    tag_r = add_rect(s, cl + Inches(0.25), ct + Inches(0.25), Inches(1.2), Inches(0.3),
                     RGBColor(int(str(color)[0:2], 16) // 5, int(str(color)[2:4], 16) // 5, int(str(color)[4:6], 16) // 5))
    add_text(s, cl + Inches(0.25), ct + Inches(0.25), Inches(1.2), Inches(0.3),
             tag, 9, color, True, PP_ALIGN.CENTER)
    add_text(s, cl + Inches(0.25), ct + Inches(0.7), cw - Inches(0.5), Inches(0.5),
             title, 15, WHITE, True)
    add_text(s, cl + Inches(0.25), ct + Inches(1.3), cw - Inches(0.5), Inches(3.0),
             body, 11, GRAY)

# Bottom insight
add_text(s, Inches(0.8), SH - Inches(1.6), Inches(11), Inches(0.8),
         "Tags are free-form SCREAMING_CASE — the SLM creates them contextually, not from a hard-coded list.\n"
         "is_essential: false → mask_and_send  |  true → prompt_user (masking would break meaning)",
         12, GRAY)

footer(s, "Detection Framework", "3 / 12")


# ════════════════════════════════════════════════════════════════════
# SLIDE 4 — Architecture
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
section_num(s, 4)
nav_breadcrumb(s, ["Title", "Problem", "Architecture", "Eval"], 2)
add_text(s, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
         "Architecture", 40, WHITE, True)
add_text(s, Inches(0.8), Inches(1.1), Inches(8), Inches(0.4),
         "How a user prompt flows through the Privacy Router pipeline", 16, GRAY)

# Flow boxes
flow_items = [
    ("USER", "Raw Prompt", "\"주민등록번호 901212-1234567로 조회해줘\""),
    ("EXTRACTOR", "SLM Detection", "Three-harm test + Critic pass"),
    ("JUDGE", "Policy Decision", "allow / mask_and_send / prompt_user"),
    ("ROUTER", "Model Selection", "Local / Cloud / User prompt"),
]
for i, (label, title, desc) in enumerate(flow_items):
    bx = Inches(0.8) + i * Inches(3.15)
    by = Inches(1.8)
    bw = Inches(2.8)
    bh = Inches(1.8)
    border = CARD if i > 0 else RGBColor(0x44, 0x44, 0x44)
    add_rect(s, bx, by, bw, bh, LIGHT_BG, border)
    add_text(s, bx + Inches(0.15), by + Inches(0.15), bw - Inches(0.3), Inches(0.25),
             label, 9, WARM, True)
    add_text(s, bx + Inches(0.15), by + Inches(0.5), bw - Inches(0.3), Inches(0.4),
             title, 14, WHITE, True)
    add_text(s, bx + Inches(0.15), by + Inches(1.0), bw - Inches(0.3), Inches(0.6),
             desc, 10, GRAY)
    if i < 3:
        ax = bx + bw + Inches(0.05)
        add_text(s, ax, by + Inches(0.7), Inches(0.3), Inches(0.3),
                 "→", 18, CARD, True)

# Bottom cards
bottom_cards = [
    ("MASKING CONTRACT", "Deterministic masking with UID placeholders. Roundtrip-safe: mask → LLM → hydrate."),
    ("TWO-PHASE EXTRACTOR", "Phase 1: SLM detects. Phase 2: Critic reviews and fills gaps. Dedup on merge."),
    ("OBSERVABILITY", "OTel traces + Prometheus metrics. Grafana dashboards. Every process() call logged."),
]
for i, (label, body) in enumerate(bottom_cards):
    cl = Inches(0.8) + i * Inches(4.1)
    ct = Inches(4.0)
    cw = Inches(3.8)
    ch = Inches(2.0)
    add_card(s, cl, ct, cw, ch, label, "", body)

footer(s, "System Design", "4 / 12")

# ════════════════════════════════════════════════════════════════════
# SLIDE 5 — Prompt Engineering v1→v2
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
section_num(s, 5)
nav_breadcrumb(s, ["Title", "Problem", "Engineering", "Eval"], 2)
add_text(s, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
         "Key Discovery: Prompt > Parameters", 40, WHITE, True)
add_text(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.4),
         "OSS models went from 0% to 100% on contextual detection — with zero additional cost", 16, GRAY)

# v1 card
add_rect(s, Inches(0.8), Inches(1.8), Inches(5.8), Inches(2.6), LIGHT_BG, RGBColor(0x33, 0x33, 0x33))
add_rect(s, Inches(0.8), Inches(1.8), Inches(0.06), Inches(2.6), RED)
add_text(s, Inches(1.1), Inches(1.9), Inches(5), Inches(0.3),
         "v1 — PATTERN MATCHING", 10, RED, True)
add_text(s, Inches(1.1), Inches(2.3), Inches(5), Inches(0.3),
         "Keyword-dependent", 16, WHITE, True)
add_text(s, Inches(1.1), Inches(2.7), Inches(5), Inches(1.2),
         'Relied on explicit cues: "주민등록번호", "결정", "채택"\n\nResult: Only Gemini 3.1 Flash Lite detected business/research secrets.\nAll open-source models: 0%.',
         12, GRAY)

# v2 card
add_rect(s, Inches(6.9), Inches(1.8), Inches(5.8), Inches(2.6), LIGHT_BG, RGBColor(0x33, 0x33, 0x33))
add_rect(s, Inches(6.9), Inches(1.8), Inches(0.06), Inches(2.6), GREEN)
add_text(s, Inches(7.2), Inches(1.9), Inches(5), Inches(0.3),
         "v2 — CONTEXTUAL REASONING", 10, GREEN, True)
add_text(s, Inches(7.2), Inches(2.3), Inches(5), Inches(0.3),
         "Question-based inference", 16, WHITE, True)
add_text(s, Inches(7.2), Inches(2.7), Inches(5), Inches(1.2),
         'Added two reasoning prompts:\n"이 문장이 내일 신문에 실리면, 경쟁사가 이득을 볼까?"\n"출판 전에 공개되면, 연구자가 피해를 볼까?"',
         12, GRAY)

# Result card
add_rect(s, Inches(0.8), Inches(4.7), Inches(11.9), Inches(1.8), RGBColor(0x26, 0x20, 0x1A), CARD)
add_text(s, Inches(1.2), Inches(4.85), Inches(2), Inches(0.3),
         "v2 RESULT", 10, WARM, True)
stats = [("1→4", "Models at 7/7"), ("$0", "Additional cost"), ("+4", "Cases improved")]
for i, (num, label) in enumerate(stats):
    sx = Inches(1.5) + i * Inches(4.0)
    add_text(s, sx, Inches(5.2), Inches(3), Inches(0.6),
             num, 36, CARD, True, PP_ALIGN.CENTER)
    add_text(s, sx, Inches(5.8), Inches(3), Inches(0.3),
             label, 12, GRAY, False, PP_ALIGN.CENTER)

footer(s, "Engineering Insight", "5 / 12")


# ════════════════════════════════════════════════════════════════════
# SLIDE 6 — Eval Cloud
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
section_num(s, 6)
nav_breadcrumb(s, ["Title", "Problem", "Engineering", "Eval"], 3)
add_text(s, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
         "Eval — Cloud Models (v2 Prompt)", 40, WHITE, True)
add_text(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.4),
         "17 Korean test cases · N=5 trials each · target_ok + context_ok metrics", 16, GRAY)

# Table — synced with HTML presentation + REPORT.md Section 7.3
rows = [
    ("Model", "Params", "Platform", "Cost/1M", "Target", "Context", "Time"),
    ("Gemma 4 26B-A4B", "26B MoE", "OpenRouter", "$0.06", "100.0%", "100.0%", "5.0s"),
    ("Gemini 3.1 Flash Lite", "—", "OpenRouter", "$0.25", "100.0%", "100.0%", "1.9s"),
    ("Gemma 4 12B", "12B", "vLLM nightly", "Self-hosted", "100.0%", "82.4%", "25.1s"),
    ("DeepSeek v4 Flash", "—", "OpenRouter", "$0.10", "100.0%", "76.5%", "—"),
    ("Gemma 4 26B-A4B (local)", "26B MoE", "vLLM BF16", "Self-hosted", "100.0%", "76.5%", "9.0s"),
    ("Gemma 4 E4B", "4B", "vLLM", "Self-hosted", "82.4%", "70.6%", "8.3s"),
    ("Gemma 4 E2B", "2B", "vLLM", "Self-hosted", "100.0%", "64.7%", "5.4s"),
    ("Ministral 3B", "3B", "OpenRouter", "$0.10", "94.1%", "52.9%", "3.2s"),
    ("Granite 4.1 8B (local)", "8B", "vLLM BF16", "Self-hosted", "100.0%", "41.2%", "19.0s"),
    ("Qwen 3.5 9B (local)", "9B", "vLLM BF16", "Self-hosted", "76.5%", "35.3%", "38.5s"),
    ("EXAONE 4.5 33B", "33B", "vLLM nightly", "Self-hosted", "58.8%", "35.3%", "12.7s"),
    ("Ministral 3B (local)", "3B", "vLLM BF16", "Self-hosted", "88.2%", "23.5%", "10.6s"),
]
col_widths = [Inches(2.4), Inches(1.0), Inches(1.4), Inches(1.2), Inches(1.0), Inches(1.0), Inches(0.9)]
start_x = Inches(0.8)
start_y = Inches(1.8)
row_h = Inches(0.36)
for ri, row in enumerate(rows):
    cx = start_x
    for ci, cell in enumerate(row):
        is_header = ri == 0
        is_high = cell in ("100.0%",)
        is_mid = cell in ("82.4%", "76.5%", "70.6%", "64.7%", "52.9%", "94.1%", "88.2%", "76.5%")
        is_low = cell in ("35.3%", "23.5%", "41.2%")
        color = WARM if is_header else (GREEN if is_high else (RGBColor(0xFF, 0x98, 0x00) if is_mid else (RED if is_low else GRAY)))
        bold = is_header or is_high or is_low
        add_text(s, cx, start_y + ri * row_h, col_widths[ci], row_h,
                 cell, 9 if is_header else 10, color, bold)
        cx += col_widths[ci]
    if ri == 0:
        add_rect(s, start_x, start_y + row_h - Inches(0.02), sum(w for w in col_widths), Inches(0.02), CARD)

# Insight
add_rect(s, Inches(0.8), Inches(6.4), Inches(11.9), Inches(0.6), RGBColor(0x26, 0x20, 0x1A), CARD)
add_text(s, Inches(1.1), Inches(6.45), Inches(11), Inches(0.5),
         "Best cost-performance: Gemma 4 26B-A4B at $0.06/1M matches Gemini Flash Lite ($0.25) — 4× cheaper.",
         12, WHITE, False)

footer(s, "Phase 1: Cloud + Local · 14 models evaluated", "6 / 12")


# ════════════════════════════════════════════════════════════════════
# SLIDE 7 — Test Cases & Evaluation
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
section_num(s, 7)
nav_breadcrumb(s, ["Title", "Problem", "Testing", "Future"], 2)
add_text(s, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
         "Evaluation Framework", 40, WHITE, True)
add_text(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.4),
         "Two distinct testing layers · 64 unit tests · 17 Korean eval scenarios", 16, GRAY)

# Left card: Software Unit Tests with coverage
add_card(s, Inches(0.8), Inches(1.8), Inches(5.8), Inches(4.5),
         "SOFTWARE UNIT TESTS", "64 Tests · Coverage by Module",
         "Verify individual components behave correctly.\n"
         "No LLM calls · No external dependencies.\n\n"
         "Coverage rates (measured via pytest-cov):\n"
         "  Masker ............ 93%\n"
         "  Extractor ......... 85%\n"
         "  Config ............ 84%\n"
         "  Judge ............. 66%\n"
         "  Router ............ 56%\n")

# Right card: LLM Agent Evaluation
add_rect(s, Inches(6.9), Inches(1.8), Inches(5.8), Inches(4.5), LIGHT_BG, RGBColor(0x33, 0x33, 0x33))
add_text(s, Inches(7.2), Inches(1.9), Inches(5), Inches(0.25),
         "LLM AGENT EVALUATION", 10, WARM, True)
add_text(s, Inches(7.2), Inches(2.2), Inches(5), Inches(0.4),
         "17 Korean Test Cases", 18, WHITE, True)
add_text(s, Inches(7.2), Inches(2.7), Inches(5), Inches(1.0),
         "Measures detection accuracy against\n"
         "real-world Korean sensitive information.\n\n"
         "Metrics: target_ok + context_ok = ok\n"
         "PII · Business · Research · Credentials · Mixed\n"
         "Models tested: 14 (cloud + local)",
         12, GRAY)

# Bottom
add_rect(s, Inches(0.8), Inches(6.5), Inches(11.9), Inches(0.6), LIGHT_BG)
add_text(s, Inches(1.1), Inches(6.55), Inches(11), Inches(0.5),
         "Coverage matrix: Every Three-Harm category × every policy action × boundary cases (empty, mixed).",
         12, GRAY)

footer(s, "Testing & Quality Assurance", "7 / 12")


# ════════════════════════════════════════════════════════════════════
# SLIDE 8 — Integration
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
section_num(s, 8)
nav_breadcrumb(s, ["Title", "Problem", "Integration", "Eval"], 2)
add_text(s, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
         "Integration", 40, WHITE, True)
add_text(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.4),
         "Single MCP tool + OpenAI-compatible proxy — works with any agent framework", 16, GRAY)

# MCP card
add_card(s, Inches(0.8), Inches(1.8), Inches(5.8), Inches(3.5),
         "MCP SERVER", "process(text, action, model)",
         "Single unified tool. action parameter controls behavior:\n\n"
         "auto     → full pipeline (detect → route → generate)\n"
         "classify → detection only (no LLM call)\n"
         "generate → mask + forward to LLM\n"
         "allow    → skip detection, direct LLM")

# OpenAI card
add_card(s, Inches(6.9), Inches(1.8), Inches(5.8), Inches(3.5),
         "OPENAI COMPATIBLE", "/v1/chat/completions",
         "Drop-in proxy for OpenAI SDK.\nAny agent using OpenAI SDK can route through\nPrivacy Router by changing the base URL.\n\n"
         "# Hermes Agent\nllm: \"http://privacy-router:8787/v1\"\n\n"
         "# OpenClaw\nmodel: \"privacy-router/model-name\"")

# Agent cards
agents = [("Hermes", "Butler agent · MCP direct"),
          ("OpenClaw", "Coding agent · OpenAI proxy"),
          ("OpenCode", "CLI agent · Custom provider")]
for i, (name, desc) in enumerate(agents):
    cl = Inches(0.8) + i * Inches(4.1)
    ct = Inches(5.6)
    cw = Inches(3.8)
    ch = Inches(1.2)
    add_rect(s, cl, ct, cw, ch, LIGHT_BG, RGBColor(0x33, 0x33, 0x33))
    add_text(s, cl + Inches(0.3), ct + Inches(0.15), cw - Inches(0.6), Inches(0.5),
             name, 20, CARD, True, PP_ALIGN.CENTER)
    add_text(s, cl + Inches(0.3), ct + Inches(0.65), cw - Inches(0.6), Inches(0.3),
             desc, 11, GRAY, False, PP_ALIGN.CENTER)

footer(s, "MCP + OpenAI Compatible API", "8 / 12")


# ════════════════════════════════════════════════════════════════════
# SLIDE 9 — 17 Test Cases
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
section_num(s, 9)
nav_breadcrumb(s, ["Title", "Problem", "Test Cases", "Future"], 2)
add_text(s, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
         "17 Korean Test Cases", 40, WHITE, True)
add_text(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.4),
         "Validated across all Three-Harm categories · Core evaluation dataset", 16, GRAY)

# Table
tc_rows = [
    ("#", "Category", "Example", "Expected"),
    ("1", "PII (context)", "주민등록번호 901212-1234567을 포함한 이메일을 작성해줘", "mask_and_send"),
    ("2", "PII (target)", "내 주민등록번호가 뭐야?", "prompt_user"),
    ("3", "Business", "TSMC 3nm 공정 채택 결정에 대한 보고서를 작성해줘", "mask_and_send"),
    ("4", "Research (ctx)", "이 새로운 Attention 대체 아이디어를 바탕으로 실험 설계를 도와줘", "mask_and_send"),
    ("5", "Research (tgt)", "새로운 Attention 대체 아이디어가 뭐야?", "prompt_user"),
    ("6", "Experiment", "이 실험 결과를 바탕으로 논문 초안을 작성해줘", "mask_and_send"),
    ("7", "Non-sensitive", "오늘 서울 날씨는 맑고 기온은 25도입니다", "allow"),
]
tc_col_widths = [Inches(0.5), Inches(1.6), Inches(7.5), Inches(2.0)]
tc_start_x = Inches(0.8)
tc_start_y = Inches(1.8)
tc_row_h = Inches(0.5)
for ri, row in enumerate(tc_rows):
    cx = tc_start_x
    is_header = ri == 0
    for ci, cell in enumerate(row):
        add_text(s, cx + Inches(0.1), tc_start_y + ri * tc_row_h,
                 tc_col_widths[ci] - Inches(0.2), tc_row_h,
                 cell, 11 if not is_header else 10,
                 WHITE if is_header else GRAY,
                 bold=is_header)
        cx += tc_col_widths[ci]
    if is_header:
        add_rect(s, tc_start_x, tc_start_y + tc_row_h - Inches(0.02),
                 sum(tc_col_widths), Inches(0.02), CARD)

add_text(s, Inches(0.8), Inches(5.7), Inches(11), Inches(0.5),
         "Cases 8–17: Mixed categories (PII + credentials + budgets + unpublished data).",
         12, GRAY)

footer(s, "Evaluation Dataset", "9 / 12")
# ════════════════════════════════════════════════════════════════════
# SLIDE 10 — Limitations
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
section_num(s, 10)
nav_breadcrumb(s, ["Title", "Problem", "Limits", "Future"], 2)
add_text(s, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
         "Honest Limitations", 40, WHITE, True)
add_text(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.4),
         "What we know we don't do well yet", 16, GRAY)

limits = [
    ("⚠", "Single-problem evaluation",
     "Each test case is evaluated independently. Multi-turn accumulation attacks are not systematically tested."),
    ("⚠", "Small test set (17 cases)",
     "17 Korean test cases cannot represent the full space of real-world prompts. Expanded, diverse test sets needed."),
    ("⚠", "Eval was single-pass per model",
     "v2 evaluation ran once (N=1). Need N=5 re-evaluation for confidence intervals."),
    ("ℹ", "Latency overhead",
     "TwoPhaseExtractor adds ~2-4s. For real-time chat, this may be noticeable. Target: sub-second."),
    ("ℹ", "Korean-only scope",
     "Tested on Korean + English. Chinese, Japanese PII patterns not validated."),
    ("ℹ", "No adversarial testing",
     "Tested with cooperative inputs. Deliberate evasion not systematically evaluated."),
]
for i, (icon, title, body) in enumerate(limits):
    col = i % 2
    row = i // 2
    cl = Inches(0.8) + col * Inches(6.3)
    ct = Inches(1.7) + row * Inches(1.75)
    cw = Inches(5.9)
    ch = Inches(1.5)
    add_rect(s, cl, ct, cw, ch, LIGHT_BG, RGBColor(0x33, 0x33, 0x33))
    add_text(s, cl + Inches(0.2), ct + Inches(0.15), Inches(0.4), Inches(0.3),
             icon, 14, WARM if icon == "⚠" else CYAN)
    add_text(s, cl + Inches(0.6), ct + Inches(0.15), cw - Inches(0.9), Inches(0.3),
             title, 14, WHITE, True)
    add_text(s, cl + Inches(0.6), ct + Inches(0.55), cw - Inches(0.9), Inches(0.8),
             body, 11, GRAY)

footer(s, "Transparency", "10 / 12")


# ════════════════════════════════════════════════════════════════════
# SLIDE 11 — Roadmap (3-Step)
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
section_num(s, 11)
nav_breadcrumb(s, ["Title", "Problem", "Limits", "Future"], 3)
add_text(s, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
         "Roadmap", 40, WHITE, True)
add_text(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.4),
         "Data Collection → Prompt Optimization → Custom Model Training", 16, GRAY)

future = [
    ("STEP 1", "Data Collection",
     "Every process() call logs detection records to PostgreSQL.\n\n"
     "Build a labeled dataset of Korean sensitive information:\n"
     "categories, confidence scores, essential decisions.\n\n"
     "Goal: Accumulate real-world training data."),
    ("STEP 2", "Prompt Optimization",
     "Continuously refine the Extractor's contextual reasoning.\n\n"
     "v1→v2 proved this: 0% → 100% at $0 cost.\n\n"
     "Goal: 99%+ accuracy on expanded test set."),
    ("STEP 3", "Custom Model Training",
     "Fine-tune or train a dedicated lightweight model\n"
     "(OpenAI Privacy Filter-style) on accumulated data.\n\n"
     "Target: 50–200ms classification latency.\n"
     "Current: 2–4s (SLM inference)."),
]
for i, (label, title, body) in enumerate(future):
    cl = Inches(0.8) + i * Inches(4.1)
    ct = Inches(1.8)
    cw = Inches(3.8)
    ch = Inches(4.5)
    fill = RGBColor(0x26, 0x20, 0x1A) if i == 2 else LIGHT_BG
    add_rect(s, cl, ct, cw, ch, fill, CARD if fill == RGBColor(0x26, 0x20, 0x1A) else RGBColor(0x33, 0x33, 0x33))
    add_text(s, cl + Inches(0.25), ct + Inches(0.2), cw - Inches(0.5), Inches(0.25),
             label, 10, WARM, True)
    add_text(s, cl + Inches(0.25), ct + Inches(0.5), cw - Inches(0.5), Inches(0.4),
             title, 18, WHITE, True)
    add_text(s, cl + Inches(0.25), ct + Inches(1.0), cw - Inches(0.5), Inches(3.0),
             body, 12, GRAY)

footer(s, "Roadmap", "11 / 12")


# ════════════════════════════════════════════════════════════════════
# SLIDE 12 — Closing
# ════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
section_num(s, 12)

# Big quote
add_text(s, Inches(0.8), Inches(1.5), Inches(7), Inches(2.5),
         "There is a lower bound on\nmodel size for contextual\nprivacy detection.",
         36, WHITE, True)
add_text(s, Inches(0.8), Inches(3.8), Inches(7), Inches(1.0),
         "Gemma 4 26B crosses that threshold —\nat $0.06 per million tokens.",
         22, GRAY)

# Summary card
card_l = SW - Inches(5.5)
card_t = Inches(1.5)
card_w = Inches(4.7)
card_h = Inches(4.5)
add_rect(s, card_l, card_t, card_w, card_h, CARD)
add_text(s, card_l + Inches(0.4), card_t + Inches(0.3), Inches(3.5), Inches(0.3),
         "SUMMARY", 10, DARK_TEXT, True)
summary = [
    ("Detection Framework", "Three-Harm Test"),
    ("Prompt Engineering", "v1→v2: 0%→100%"),
    ("Best Model", "Gemma 4 26B · $0.06/1M"),
    ("Integration", "MCP + OpenAI Proxy"),
    ("Test Coverage", "47 unit + 17 eval"),
    ("Roadmap", "Data → Prompt → Model"),
]
for i, (label, value) in enumerate(summary):
    sy = card_t + Inches(0.8) + i * Inches(0.6)
    add_text(s, card_l + Inches(0.4), sy, Inches(2.2), Inches(0.3),
             label, 13, DARK_TEXT)
    add_text(s, card_l + Inches(2.6), sy, Inches(1.8), Inches(0.3),
             value, 11, DARK_TEXT, False, PP_ALIGN.RIGHT)

footer(s, "Privacy Router · GenAI & Blockchain · June 2026", "DH. Kim & M. Saadati")

# ── Save ────────────────────────────────────────────────────────────
out = "PrivacyRouter_Presentation.pptx"
prs.save(out)
print(f"✓ Saved {out} ({len(prs.slides)} slides)")
